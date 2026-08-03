#!/usr/bin/env python3
"""Презентация «Legal Design — стандарт документов 2026».

Вывод — на Рабочий стол. Дизайн: минимализм, крупная типографика, воздух,
один акцентный цвет, никакого клипарта. Все изображения — реальные рендеры
собранных образцов .docx, не мокапы.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
DECK_IMG = HERE / "samples" / "_deck"
OUT = Path.home() / "Desktop" / "Legal Design — стандарт документов 2026.pptx"

W, H = Inches(13.333), Inches(7.5)

PAPER = RGBColor(0xFB, 0xFA, 0xF8)
INK = RGBColor(0x0F, 0x11, 0x13)
GREY = RGBColor(0x70, 0x73, 0x78)
HAIR = RGBColor(0xDC, 0xD7, 0xCE)
ACCENT = RGBColor(0xA6, 0x3A, 0x2A)
DARK = RGBColor(0x14, 0x16, 0x18)

DISPLAY = "Avenir Next"
BODY = "Helvetica Neue"

M = Inches(0.95)          # боковое поле
CW = W - 2 * M            # ширина полосы


def deck():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def slide(prs, bg=PAPER):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def text(s, txt, left, top, width, height, size, font=BODY, color=INK,
         bold=False, align=PP_ALIGN.LEFT, spacing=1.0, space_after=0,
         anchor=MSO_ANCHOR.TOP, caps=False):
    box = s.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = txt.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        para.space_after = Pt(space_after)
        run = para.add_run()
        run.text = line.upper() if caps else line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def rule(s, left, top, width, color=HAIR, thick=Pt(1)):
    ln = s.shapes.add_connector(1, left, top, left + width, top)
    ln.line.color.rgb = color
    ln.line.width = thick
    return ln


def kicker(s, txt, top=Inches(0.62), color=ACCENT):
    text(s, txt, M, top, CW, Inches(0.3), 11.5, DISPLAY, color,
         bold=True, caps=True, spacing=1.0)
    rule(s, M, top + Inches(0.36), CW)


def picture(s, name, left, top, width=None, height=None):
    path = DECK_IMG / f"{name}.png"
    if width is not None:
        return s.shapes.add_picture(str(path), left, top, width=width)
    return s.shapes.add_picture(str(path), left, top, height=height)


def caption(s, txt, left, top, width, align=PP_ALIGN.LEFT):
    text(s, txt, left, top, width, Inches(0.4), 10.5, BODY, GREY, align=align,
         spacing=1.15)


def frame(s, shape, pad=Emu(0)):
    """Тонкая рамка вокруг изображения — отделяет скан документа от полосы."""
    box = s.shapes.add_shape(1, shape.left - pad, shape.top - pad,
                             shape.width + 2 * pad, shape.height + 2 * pad)
    box.fill.background()
    box.line.color.rgb = HAIR
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    return box


# ------------------------------------------------------------------ слайды
def s_title(prs):
    s = slide(prs, DARK)
    text(s, "Legal design", M, Inches(2.15), CW, Inches(1.3), 74, DISPLAY,
         PAPER, bold=True, spacing=0.92)
    text(s, "стандарт документов 2026", M, Inches(3.32), CW, Inches(1.0), 40,
         DISPLAY, RGBColor(0x9A, 0x9D, 0xA2), spacing=1.0)
    rule(s, M, Inches(4.55), Inches(2.2), ACCENT, Pt(2.5))
    text(s, "Зарубин и партнеры · 03.08.2026", M, Inches(4.85), CW,
         Inches(0.5), 14, BODY, RGBColor(0x8A, 0x8D, 0x92), spacing=1.3)
    text(s, "Три независимых исследовательских потока · перекрестная критика ·\n"
            "два собранных образца .docx · протокол испытаний",
         M, Inches(6.05), Inches(9.0), Inches(1.0), 12.5, BODY,
         RGBColor(0x6E, 0x71, 0x76), spacing=1.35)
    return s


def s_question(prs):
    s = slide(prs)
    kicker(s, "Зачем это все")
    text(s, "Как теперь выглядит\nнаш документ\nи почему именно так", M,
         Inches(1.55), Inches(8.4), Inches(3.2), 46, DISPLAY, INK, bold=True,
         spacing=1.05)
    rule(s, M, Inches(4.75), Inches(1.6), ACCENT, Pt(2.5))
    cols = [
        ("Что доказано", "Замеры есть, но их мало и они\nне про то, о чем обычно говорят."),
        ("Что можно в РФ", "Суд ограничивает содержание\nи канал подачи, а не верстку."),
        ("Что собрано", "Два образца .docx, проверенных\nрендером, а не обещанием."),
    ]
    x = M
    cw = Inches(3.55)
    for head, body in cols:
        text(s, head, x, Inches(5.25), cw, Inches(0.4), 15, DISPLAY, INK,
             bold=True)
        text(s, body, x, Inches(5.72), cw, Inches(1.2), 12.5, BODY, GREY,
             spacing=1.35)
        x += cw + Inches(0.35)
    return s


def s_statement(prs):
    s = slide(prs, DARK)
    text(s, "Российский суд почти не ограничивает\nвизуальную форму документа стороны.",
         M, Inches(2.05), CW, Inches(2.0), 36, DISPLAY, PAPER, bold=True,
         spacing=1.18)
    text(s, "Он жестко ограничивает содержание, подпись и канал подачи.\n"
            "Ни одного акта о возврате по мотиву оформления не найдено ни в одном из трех потоков.",
         M, Inches(4.25), CW, Inches(1.2), 17, BODY,
         RGBColor(0x9A, 0x9D, 0xA2), spacing=1.45)
    rule(s, M, Inches(5.75), Inches(1.6), ACCENT, Pt(2.5))
    text(s, "Препятствие не в праве. Препятствие в привычке.",
         M, Inches(6.05), CW, Inches(0.5), 14, BODY,
         RGBColor(0x70, 0x73, 0x78), spacing=1.3)
    return s


def s_proven(prs):
    s = slide(prs)
    kicker(s, "Что действительно замерено")
    text(s, "Четыре цифры, на которые можно ссылаться", M, Inches(1.2), CW,
         Inches(0.6), 30, DISPLAY, INK, bold=True)
    items = [
        ("23 → 70 %", "понимание цели судебной формы после\nупрощения языка",
         "Mindlin, Judicial Council of California, 2005\n60 респондентов, значимо при 99 %"),
        ("66 / 34", "судьи предпочли plain English варианту\nна legalese",
         "Flammer, J. Legal Writing Institute, 2010\n149 выраженных предпочтений"),
        ("897 / 1080 с", "время понимания договора: диаграмма\nпротив сплошного текста",
         "Passera, Kankaanranta, Louhiala-Salminen,\nIEEE TPC, 2017 · p = 0,004"),
        ("1–2 из 1000", "покупателей вообще открывают\nлицензионное соглашение",
         "Bakos, Marotta-Wurgler, Trossen,\nJ. Legal Studies, 2014"),
    ]
    x = M
    cw = Inches(2.62)
    for big, mid, src in items:
        text(s, big, x, Inches(2.35), cw, Inches(0.75), 33, DISPLAY, ACCENT,
             bold=True)
        text(s, mid, x, Inches(3.15), cw, Inches(1.1), 13, BODY, INK,
             spacing=1.32)
        rule(s, x, Inches(4.55), Inches(1.0))
        text(s, src, x, Inches(4.75), cw, Inches(1.0), 10, BODY, GREY,
             spacing=1.3)
        x += cw + Inches(0.28)
    text(s, "Последняя цифра — против всего движения: если документ не открывают, "
            "внутри него дизайн уже не работает.",
         M, Inches(6.35), CW, Inches(0.5), 13, BODY, GREY, spacing=1.35)
    return s


def s_not_proven(prs):
    s = slide(prs)
    kicker(s, "Чего не доказал никто")
    text(s, "Границы честного разговора с доверителем", M, Inches(1.2), CW,
         Inches(0.6), 30, DISPLAY, INK, bold=True)
    rows = [
        ("Диаграмма в иске меняет решение судьи",
         "Эксперимента не существует. Ни в одном отчете, ни в одной рецензии."),
        ("Читаемость выигрывает дело",
         "Два замера противоречат друг другу: 31 → 69 % в федеральном суде США "
         "(Spencer, Feldman, 2018) и полное отсутствие связи на 882 брифах (Long, Christensen)."),
        ("Слоеный договор улучшает понимание",
         "Замер против: слоеная версия отработала не лучше полного текста, "
         "до глубокого слоя почти никто не доходит (Kelley, Cranor, CHI 2010, 764 участника)."),
        ("«Вывод первый» сам по себе работает",
         "Перестановка вывода в начало без упрощения языка понимания не улучшает."),
        ("Иконки, цвет, комикс-форма",
         "Самостоятельного замера эффекта нет ни у одной. Стандарта иконок нет даже в ЕС: "
         "делегированный акт по ст. 12(8) GDPR не принят девять лет."),
    ]
    y = Inches(2.15)
    for head, body in rows:
        text(s, head, M, y, Inches(4.2), Inches(0.5), 14.5, DISPLAY, INK,
             bold=True, spacing=1.2)
        text(s, body, M + Inches(4.5), y, CW - Inches(4.5), Inches(0.8), 12.5,
             BODY, GREY, spacing=1.35)
        y += Inches(0.94)
        rule(s, M, y - Inches(0.22), CW)
    return s


def s_before_after(prs):
    s = slide(prs)
    kicker(s, "До и после · одна фактура, одни нормы")
    half = Emu(int((CW - Inches(0.5)) / 2))
    text(s, "Как обычно", M, Inches(1.28), half, Inches(0.4), 16, DISPLAY,
         GREY, bold=True)
    text(s, "По стандарту", M + half + Inches(0.5), Inches(1.28), half,
         Inches(0.4), 16, DISPLAY, ACCENT, bold=True)
    pic1 = picture(s, "do_page1", M, Inches(1.85), height=Inches(4.75))
    pic1.left = M + int((half - pic1.width) / 2)
    frame(s, pic1)
    pic2 = picture(s, "after_page1", M + half + Inches(0.5), Inches(1.85),
                   height=Inches(4.75))
    pic2.left = M + half + Inches(0.5) + int((half - pic2.width) / 2)
    frame(s, pic2)
    caption(s, "Три абзаца по 200 слов. Расчет прозой.\nХронология внутри предложения.",
            M, Inches(6.75), half)
    caption(s, "Суть спора в рамке. Оглавление. Нумерованные абзацы.\n"
               "Расчет таблицей. Таймлайн и схема сторон.",
            M + half + Inches(0.5), Inches(6.75), half)
    return s


def s_calc(prs):
    s = slide(prs)
    kicker(s, "До и после · расчет процентов")
    pic1 = picture(s, "do_calc", M, Inches(1.35), width=Inches(5.7))
    frame(s, pic1)
    caption(s, "Прозой: 118 слов в одном предложении. Проверить арифметику "
               "невозможно, не выписав числа на бумагу.",
            M, Inches(1.35) + pic1.height + Inches(0.25), Inches(5.7))
    pic2 = picture(s, "isk_calc", M + Inches(6.2), Inches(1.35),
                   width=Inches(5.5))
    frame(s, pic2)
    caption(s, "Таблицей: периоды, ставка, дни, итог. Судья сверяет "
               "глазами, не пересчитывая.",
            M + Inches(6.2), Inches(1.35) + pic2.height + Inches(0.25),
            Inches(5.5))
    return s


def s_element(prs, kick, head, img, note, img_w=Inches(8.1)):
    s = slide(prs)
    kicker(s, kick)
    text(s, head, M, Inches(1.2), Inches(11.4), Inches(0.6), 28, DISPLAY, INK,
         bold=True)
    pic = picture(s, img, M, Inches(2.15), width=img_w)
    if pic.top + pic.height > Inches(6.85):
        pic.height = Inches(6.85) - pic.top
        pic.width = int(pic.height * 1.0)
    frame(s, pic)
    text(s, note, M + img_w + Inches(0.5), Inches(2.15),
         CW - img_w - Inches(0.5), Inches(3.0), 13, BODY, GREY, spacing=1.42)
    return s


def s_patterns(prs):
    s = slide(prs)
    kicker(s, "Свод паттернов")
    text(s, "Что берем, при каком условии и как собирается", M, Inches(1.15),
         CW, Inches(0.5), 28, DISPLAY, INK, bold=True)
    rows = [
        ("Паттерн", "Доказательства", "В РФ", "Сборка в .docx"),
        ("Простой язык", "замер", "свободно", "текст"),
        ("Таблица вместо текста", "замер", "свободно", "нативная таблица — самый надежный элемент"),
        ("Резюме в начале", "норма за рубежом", "с оговоркой", "таблица 1×1 с рамкой"),
        ("Нумерация абзацев", "норма в 4 системах", "свободно", "висячий отступ"),
        ("Оглавление", "норма в США", "свободно", "статический список; поле TOC не годится"),
        ("Таймлайн, схема сторон", "нет замера", "с оговоркой", "PNG 300 dpi, черно-белый"),
        ("Диаграмма в договоре", "замер", "вне процесса", "PNG"),
        ("Двухслойный договор", "замер против", "только с пунктом о приоритете", "два слоя в одном файле"),
        ("Цвет, иконки", "нет замера", "неприменим в суде", "хрупко, теряется в ч/б"),
        ("Интерактив", "—", "запрещен", "не собирать"),
    ]
    top = Inches(2.05)
    widths = [Inches(3.05), Inches(2.15), Inches(2.35), Inches(3.9)]
    tbl_w = sum(widths, Emu(0))
    shape = s.shapes.add_table(len(rows), 4, M, top, tbl_w,
                               Inches(0.42) * len(rows))
    table = shape.table
    for i, wdt in enumerate(widths):
        table.columns[i].width = wdt
    for r, row in enumerate(rows):
        table.rows[r].height = Inches(0.40)
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            cell.margin_left = Inches(0.09)
            cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER if r else RGBColor(0xEE, 0xEB, 0xE5)
            para = cell.text_frame.paragraphs[0]
            para.line_spacing = 1.0
            run = para.add_run()
            run.text = val
            run.font.name = DISPLAY if r == 0 else BODY
            run.font.size = Pt(11 if r == 0 else 10.5)
            run.font.bold = r == 0
            run.font.color.rgb = INK if r == 0 else (
                ACCENT if c == 2 and val in ("запрещен", "неприменим в суде") else GREY)
    text(s, "Полная таблица с юрисдикциями и нормами — knowledge/legal-design/legal-design-standard.md",
         M, Inches(6.85), CW, Inches(0.4), 10.5, BODY, GREY)
    return s


def s_ru_limits(prs):
    s = slide(prs, DARK)
    kicker(s, "Ограничения российского суда", color=ACCENT)
    text(s, "Что ломает половину модного дизайна", M, Inches(1.2), CW,
         Inches(0.6), 30, DISPLAY, PAPER, bold=True)
    left_items = [
        ("Запрещено технически",
         "Интерактив, сворачиваемые блоки, мультимедиа, встроенные сценарии, "
         "защита от копирования.\nПриказы Судебного департамента № 251 (суды общей "
         "юрисдикции) и № 252 (арбитражные суды). Путать их нельзя — разные "
         "требования к подписи и маршруту."),
        ("Единственная норма про вид документа",
         "«Административное исковое заявление подается в суд в письменной форме "
         "в разборчивом виде» — ч. 1 ст. 125 КАС РФ. Ни одного дела ее применения "
         "не найдено: норма спящая."),
    ]
    right_items = [
        ("Физическое ограничение",
         "Правый нижний угол первого листа обязан быть свободен от текста — там "
         "ставят регистрационный штамп (п. 3.1.5 Инструкции, утв. Постановлением "
         "Пленума ВАС РФ от 25.12.2013 № 100). Основанием возврата не является, "
         "но угол защищаем."),
        ("Что режет на самом деле",
         "Содержание по перечню кодекса, подпись и полномочия, канал подачи. "
         "Реальные возвраты — про них, а не про оформление."),
    ]
    colw = Emu(int((CW - Inches(0.7)) / 2))
    for col, items in ((M, left_items), (M + colw + Inches(0.7), right_items)):
        y = Inches(2.1)
        for head, body in items:
            text(s, head, col, y, colw, Inches(0.4), 15, DISPLAY, ACCENT,
                 bold=True)
            text(s, body, col, y + Inches(0.42), colw, Inches(1.7), 12.5, BODY,
                 RGBColor(0xA8, 0xAB, 0xB0), spacing=1.42)
            y += Inches(2.35)
    return s


def s_docx_limits(prs):
    s = slide(prs)
    kicker(s, "Протокол испытаний .docx")
    text(s, "Три приема не пережили проверку рендером", M, Inches(1.2), CW,
         Inches(0.6), 30, DISPLAY, INK, bold=True)
    rows = [
        ("Автособираемое оглавление",
         "Поле TOC вычисляет только Word и только по команде обновления полей. "
         "В любом другом просмотрщике на его месте пустая строка.",
         "Статический список разделов"),
        ("Рамка вокруг абзаца",
         "XML схемно верен, все четыре стороны на месте — рендерер показал одну "
         "нижнюю границу и проглотил заливку. На группе абзацев ломается при "
         "разрыве страницы.",
         "Таблица 1×1 с рамкой и заливкой"),
        ("Внутренняя гиперссылка",
         "Рендерер красит ее синим с подчеркиванием поверх прямого форматирования. "
         "Явные color и underline в схемно верном порядке не помогли. "
         "Подчеркивание нарушает наш эталон.",
         "Оглавление обычным текстом; ссылки — только в договоре"),
    ]
    y = Inches(2.1)
    for head, why, fix in rows:
        text(s, head, M, y, Inches(3.1), Inches(0.5), 14.5, DISPLAY, INK,
             bold=True, spacing=1.2)
        text(s, why, M + Inches(3.35), y, Inches(5.3), Inches(1.1), 12, BODY,
             GREY, spacing=1.35)
        text(s, fix, M + Inches(8.9), y, Inches(2.5), Inches(1.1), 12, BODY,
             ACCENT, spacing=1.35)
        y += Inches(1.42)
        rule(s, M, y - Inches(0.3), CW)
    text(s, "Урок шире зондов: OOXML требует строгого порядка дочерних элементов. "
            "Элемент, дописанный в конец, схема считает недопустимым, и рендерер "
            "молча его игнорирует — ни ошибки, ни предупреждения.",
         M, Inches(6.55), CW, Inches(0.7), 12.5, BODY, INK, spacing=1.4)
    return s


def s_checklist(prs):
    s = slide(prs)
    kicker(s, "Чек-лист по умолчанию")
    colw = Emu(int((CW - Inches(0.7)) / 2))
    text(s, "Процессуальный документ", M, Inches(1.2), colw, Inches(0.5), 21,
         DISPLAY, INK, bold=True)
    text(s, "Договор, оферта, политика, претензия", M + colw + Inches(0.7),
         Inches(1.2), colw, Inches(0.5), 21, DISPLAY, INK, bold=True)
    left = [
        "Первый абзац — вывод, простым языком",
        "Резюме-блок: суть · требование · цена · ключевой документ",
        "Оглавление от 6 страниц, обычным текстом",
        "Сквозная нумерация абзацев от 4 страниц",
        "Источник в тексте у каждого правового утверждения",
        "Нормы дословно через cite.py либо не цитировать",
        "Расчет только таблицей, итог жирным",
        "Таймлайн от пяти событий, схема сторон от трех лиц",
        "Раздел «доводы оппонента» — в апелляции и кассации",
        "Ни одного смысла только в цвете или только в картинке",
        "Правый нижний угол первого листа свободен",
        "Прогнан scan_legal.sh, проверен рендером глазами",
    ]
    right = [
        "Краткий слой: кто · что · сколько · когда · что если",
        "Пункт о приоритете полного текста над кратким слоем",
        "Карта обязательств таблицей: сторона · срок · последствие",
        "Таймлайн сроков от трех этапов",
        "Обременительные условия отдельным выделенным разделом",
        "Чеклист приемки при поэтапном исполнении",
        "Гарнитура общедоступная — Arial, Verdana, Georgia",
        "Термин вводится один раз, дальше единообразно",
        "Есть линейная печатная версия, если пойдет в суд",
    ]
    for col, items in ((M, left), (M + colw + Inches(0.7), right)):
        y = Inches(1.95)
        for item in items:
            text(s, "—", col, y, Inches(0.25), Inches(0.3), 12, BODY, ACCENT)
            text(s, item, col + Inches(0.32), y, colw - Inches(0.32),
                 Inches(0.35), 12.5, BODY, INK, spacing=1.2)
            y += Inches(0.415)
    return s


def s_white(prs):
    s = slide(prs)
    kicker(s, "Белые пятна")
    text(s, "Что проверяем на живом деле", M, Inches(1.2), CW, Inches(0.6), 30,
         DISPLAY, INK, bold=True)
    items = [
        ("Черно-белая печать в судах",
         "Ни нормы, ни акта. От этого зависит допустимость цвета и графики. "
         "Подать документ с серой заливкой и запросить материалы дела."),
        ("Поведение приемов в самом Word",
         "Все выводы получены на рендерере QuickLook. Рамка абзаца и внутренняя "
         "ссылка в Word, вероятно, работают иначе."),
        ("ГК РФ для договорного класса",
         "Ни один из трех потоков не разобрал ст. 420, 421, 428, 431, 432, 434 "
         "применительно к форме и толкованию. Самый крупный пробел исследования."),
        ("Измеримость для русского языка",
         "Формулы Flesch откалиброваны на английском. Без адаптированного индекса "
         "рекомендация «мерить читаемость» неисполнима."),
        ("Своя эмпирика",
         "Ни одного замера на нашем материале. Пока его нет, все, что мы говорим "
         "доверителю о пользе дизайна, — заимствованные чужие цифры."),
    ]
    y = Inches(2.15)
    for head, body in items:
        text(s, head, M, y, Inches(3.7), Inches(0.5), 14, DISPLAY, INK,
             bold=True, spacing=1.2)
        text(s, body, M + Inches(4.0), y, CW - Inches(4.0), Inches(0.8), 12.5,
             BODY, GREY, spacing=1.35)
        y += Inches(0.92)
        rule(s, M, y - Inches(0.2), CW)
    return s


def s_final(prs):
    s = slide(prs, DARK)
    text(s, "Наш документ теперь\nначинается с вывода,\nсчитается таблицей\nи ссылается на пункт,\nа не на страницу.",
         M, Inches(1.55), Inches(9.5), Inches(4.2), 34, DISPLAY, PAPER,
         bold=True, spacing=1.22)
    rule(s, M, Inches(6.05), Inches(1.6), ACCENT, Pt(2.5))
    text(s, "Все остальное — цвет, иконки, слои, комиксы — осталось за бортом "
            "не потому, что запрещено, а потому, что никем не измерено.",
         M, Inches(6.35), Inches(10.8), Inches(0.8), 14, BODY,
         RGBColor(0x8A, 0x8D, 0x92), spacing=1.4)
    return s


def build():
    prs = deck()
    s_title(prs)
    s_question(prs)
    s_statement(prs)
    s_proven(prs)
    s_not_proven(prs)
    s_before_after(prs)
    s_calc(prs)
    s_element(prs, "Что появилось · 1", "Суть спора до того, как начнется спор",
              "isk_summary",
              "Судья читает иск 10-15 минут. Блок в рамке отдает сразу четыре вещи: "
              "что произошло, чего просим, сколько и каким документом подтверждается.\n\n"
              "Требование права: блок дополняет обязательные сведения, а не заменяет их. "
              "Кодекс требует, чтобы сведения были указаны в тексте.\n\n"
              "Собран таблицей 1×1 — рамка абзаца в .docx ненадежна.",
              img_w=Inches(7.6))
    s_element(prs, "Что появилось · 2", "Навигация и нумерация абзацев",
              "isk_toc_num",
              "Оглавление — от шести страниц. Обычным текстом: поле TOC вне Word "
              "не вычисляется, а гиперссылка красится синим с подчеркиванием.\n\n"
              "Сквозная нумерация абзацев — от четырех страниц. Дает точную ссылку "
              "«пункт 14» вместо «страница 6», которая поедет при любой правке.\n\n"
              "Норма в четырех системах: США, Англия, Верховный суд Великобритании, "
              "суды DIFC. Замера эффекта нет ни у кого — берем по цене внедрения.",
              img_w=Inches(7.6))
    s_element(prs, "Что появилось · 3", "Хронология и связи сторон",
              "isk_timeline",
              "Таймлайн — от пяти событий. Схема сторон — от трех лиц.\n\n"
              "Только черно-белые: смысл, живущий в цвете, не переживает серый скан "
              "и черно-белую печать.\n\n"
              "Иллюстрирует, не заменяет: те же даты обязаны стоять в тексте. "
              "Расхождение схемы и текста — готовый довод оппонента.",
              img_w=Inches(8.2))
    s_element(prs, "Договор · 1", "Главное на одной странице",
              "dog_layer",
              "Кто · что · сколько · когда · что будет, если нарушить.\n\n"
              "Обязательная страховка — пункт о приоритете полного текста над "
              "кратким слоем. Без него расхождение слоев становится предметом спора, "
              "а не удобством.\n\n"
              "Замер по слоям, к слову, отрицательный: до глубокого слоя почти никто "
              "не доходит. Берем как удобство первого контакта, не как средство "
              "добиться понимания.",
              img_w=Inches(7.3))
    s_element(prs, "Договор · 2", "Карта обязательств вместо прозы",
              "dog_oblig",
              "Сторона · обязательство · срок · последствие нарушения.\n\n"
              "Единственный паттерн с замером в обе стороны: стандартизированная "
              "табличная подача выигрывает у сплошного текста по пониманию, скорости "
              "и сравнимости (Kelley, Cranor, CHI 2010, 764 участника).\n\n"
              "Нативная таблица Word — самый устойчивый элемент из всех проверенных.",
              img_w=Inches(7.3))
    s_element(prs, "Договор · 3", "Условия, требующие отдельного внимания",
              "dog_attention",
              "Ограничение ответственности, право на замену, подсудность — "
              "выделены и собраны в один блок.\n\n"
              "В КНР это прямая обязанность: ст. 496 ГК КНР требует «разумным способом "
              "привлечь внимание» к обременительным условиям, иначе условие не станет "
              "частью договора.\n\n"
              "В РФ ближайшая опора — п. 2 ст. 8 Закона о защите прав потребителей: "
              "информация доводится «в наглядной и доступной форме». Ни один из трех "
              "потоков эту норму не нашел — ее нашел совет.",
              img_w=Inches(7.3))
    s_patterns(prs)
    s_ru_limits(prs)
    s_docx_limits(prs)
    s_checklist(prs)
    s_white(prs)
    s_final(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(OUT)
    return OUT


if __name__ == "__main__":
    build()
